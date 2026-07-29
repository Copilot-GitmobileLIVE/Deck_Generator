"""
builder.py — PPTBuilder

Responsibility:
    Orchestrate the creation of the PPTX file by iterating over all slides
    in order, delegating text rendering to SlideRenderer and image insertion
    to ImageRenderer.

Layer ordering matters in PPTX:
    python-pptx renders shapes in the order they are added to a slide.
    The first shape added is at the back (z-index 0).  This means:

    For TITLE / CLOSING slides (full-bleed image with overlaid text):
        1. ImageRenderer adds the image FIRST  → it sits behind everything
        2. SlideRenderer adds text boxes ON TOP of the image

    For CONTENT slides (side-by-side layout):
        1. SlideRenderer adds text boxes on the left
        2. ImageRenderer adds the image on the right
        (z-order does not matter here since they don’t overlap)

PPTX canvas:
    Width  = 13.33 inches  (standard widescreen 16:9)
    Height =  7.50 inches
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches

from deck_generator.models import LayoutSpec, SlideSpec
from .image_renderer import ImageRenderer
from .slide_renderer import SlideRenderer

logger = logging.getLogger("deck_generator.pptx_builder")


class PPTBuilder:
    """Top-level builder that wires SlideRenderer and ImageRenderer together.

    Usage::

        builder = PPTBuilder(output_path="output/deck.pptx")
        builder.build(slides=slides, layouts=layouts, images=images)
    """

    def __init__(self, output_path: str) -> None:
        self.output_path = output_path
        # Create a blank Presentation object — python-pptx's entry point.
        self._prs = Presentation()
        # Override the default slide dimensions to 16:9 widescreen.
        # MUST be set before adding any slides, otherwise existing slides
        # keep the default 10x7.5 dimensions.
        self._prs.slide_width = Inches(13.33)
        self._prs.slide_height = Inches(7.5)
        # Instantiate the two rendering helpers (stateless, can be reused).
        self._slide_renderer = SlideRenderer()
        self._image_renderer = ImageRenderer()

    @property
    def _blank_layout(self):
        """Return the 'Blank' slide layout — no pre-placed title/content placeholders.

        Using a blank layout is essential because we position every text box and
        image manually via Inches() coordinates.  A layout with built-in
        placeholders would interfere with our custom positioning.
        """
        layouts = self._prs.slide_layouts
        # Search for a layout explicitly named "Blank" first.
        for layout in layouts:
            if layout.name.lower() == "blank":
                return layout
        # If none found (non-standard PPTX theme), fall back to the last layout
        # which is conventionally blank in most themes.
        return layouts[-1]

    def build(
        self,
        slides: List[SlideSpec],
        layouts: Dict[int, LayoutSpec],
        images: Dict[int, str],
    ) -> str:
        """Render all slides and save the PPTX.

        Args:
            slides:  Ordered list of :class:`SlideSpec` objects.
            layouts: Map from ``slide_number`` → :class:`LayoutSpec`.
            images:  Map from ``slide_number`` → local image file path.

        Returns:
            Absolute path to the saved PPTX file.
        """
        # Sort by slide_number so the order in the PPTX matches the narrative
        # even if the list arrived in a different order.
        for spec in sorted(slides, key=lambda s: s.slide_number):
            layout = layouts.get(spec.slide_number)
            if layout is None:
                # This should not happen if LayoutAgent ran successfully,
                # but we skip gracefully rather than crashing the whole deck.
                logger.warning(
                    "PPTBuilder: no layout for slide %d — skipping", spec.slide_number
                )
                continue

            # Add a new blank slide to the presentation.
            pptx_slide = self._prs.slides.add_slide(self._blank_layout)
            image_path: Optional[str] = images.get(spec.slide_number)

            # Layer ordering: title/closing slides need the image BEHIND the text.
            if spec.slide_type.value in ("title", "closing"):
                # Image first → it sits at z-index 0 (background)
                self._image_renderer.render(pptx_slide, layout, image_path)
                # Text second → it renders on top of the image
                self._slide_renderer.render(pptx_slide, spec, layout)
            else:
                # Content slides: text on the left, image on the right.
                # Render text first (though z-order doesn’t matter here
                # since the two regions don’t overlap).
                self._slide_renderer.render(pptx_slide, spec, layout)
                self._image_renderer.render(pptx_slide, layout, image_path)

            logger.debug("PPTBuilder: rendered slide %02d — %s", spec.slide_number, spec.title)

        # Ensure the parent directory exists, then save.
        Path(self.output_path).parent.mkdir(parents=True, exist_ok=True)
        self._prs.save(self.output_path)
        logger.info("PPTBuilder: saved → %s", self.output_path)
        return self.output_path
