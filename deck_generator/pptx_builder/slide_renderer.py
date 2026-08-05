"""
slide_renderer.py — SlideRenderer

Responsibility:
    Add every text and shape element to a single python-pptx Slide object
    using the positioning values from a LayoutSpec.

All measurements use python-pptx's Inches() helper which converts inches to
English Metric Units (EMU) internally.  Font sizes use Pt() (points).

ML arteka brand header lockup (content and agenda slides only):
    The header lockup runs top-to-bottom at fixed y positions:
        y 0.26" — Eyebrow   ALL CAPS label (Indigo Grey on light, Orange on dark)
        y 0.56" — Tick rule  short orange rectangle, 0.77" wide
        y 0.85" — Title      Bold, 18-20pt, sentence case
        y 1.55" — Intro line one plain framing sentence, regular weight

Rendering order per slide (shapes added in z-order, back to front):
    1. Background fill  — applied to slide.background (always behind all shapes)
    2. Eyebrow text box — brand header lockup; skipped when show_brand_header=False
    3. Tick rule shape  — orange 0.77" accent rule; skipped with header
    4. Title text box
    5. Intro line       — skipped when show_brand_header=False
    6. Key message / subtitle text box
    7. Bullet points text box  (one text box; one paragraph per bullet)
    8. Takeaway bar    — full-width navy bar + centered text; content slides only
    9. Slide number    — bottom-right; omitted on title and closing slides
   10. Speaker notes   — written to the notes pane, not visible on the slide
"""
from __future__ import annotations

import logging
from typing import Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn as _xml_qn
from pptx.slide import Slide
from pptx.util import Inches, Pt

from deck_generator.models import LayoutSpec, SlideSpec, SlideType

# Dense consulting layout zone coordinates (13.33 × 7.5" widescreen canvas, all in inches)
_DC_LEFT_X  = 0.50   # left zone x start
_DC_LEFT_W  = 7.65   # left zone width
_DC_DIV_X   = 8.40   # vertical divider x
_DC_RIGHT_X = 8.58   # right panel x start
_DC_RIGHT_W = 4.25   # right panel width
_DC_TOP     = 2.05   # both zones y start (below intro line)
_DC_BOTTOM  = 6.28   # both zones y end (just above takeaway bar)

logger = logging.getLogger("deck_generator.slide_renderer")


def _rgb(hex_color: str) -> RGBColor:
    """Convert a '#RRGGBB' hex string to a python-pptx RGBColor object.

    python-pptx requires RGBColor(r, g, b) with integer channel values 0–255.
    We store colours as hex strings in LayoutSpec so they are human-readable
    and JSON-serialisable; this helper converts them at render time.

    Example:
        _rgb("#0057B8")  →  RGBColor(0, 87, 184)
    """
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class SlideRenderer:
    """Renders text content onto a slide using LayoutSpec positioning rules.

    All public-facing logic is in render().  The private _add_* methods each
    handle one visual layer and are separated for readability and testability.
    """

    def _set_background(self, slide: Slide, layout: LayoutSpec) -> None:
        """Fill the entire slide background with the layout's background colour."""
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(layout.background_color)

    def _add_eyebrow(self, slide: Slide, spec: SlideSpec, layout: LayoutSpec) -> None:
        """Add the ALL CAPS eyebrow label above the title (brand header lockup).

        Light slides: Indigo Grey #434E80 (AA-safe on Warm Peach).
        Dark slides: Orange #E9590C (AA-safe at any size on navy).
        """
        if not layout.show_brand_header:
            return
        text = spec.eyebrow.upper() if spec.eyebrow else spec.title[:30].upper()
        box = slide.shapes.add_textbox(
            Inches(0.5), Inches(layout.eyebrow_top_inches),
            Inches(9.0), Inches(0.28),
        )
        tf = box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.bold = True
        run.font.size = Pt(9)
        run.font.color.rgb = _rgb(layout.eyebrow_color)
        run.font.name = layout.font_family
        # letter spacing: charSpacing equivalent via XML is not needed; standard tracking

    def _add_tick_rule(self, slide: Slide, layout: LayoutSpec) -> None:
        """Add a short orange tick rule below the eyebrow (~0.77\" wide, 2pt tall)."""
        if not layout.show_brand_header:
            return
        rule = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.5), Inches(0.56),
            Inches(0.77), Inches(0.03),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = _rgb("#E9590C")  # Always orange, per brand spec
        rule.line.width = 0

    def _add_title(self, slide: Slide, spec: SlideSpec, layout: LayoutSpec) -> None:
        """Add the main title text box at the position defined in the layout."""
        box = slide.shapes.add_textbox(
            Inches(layout.title_left_inches),
            Inches(layout.title_top_inches),
            Inches(layout.title_width_inches),
            Inches(layout.title_height_inches),
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = spec.title
        run.font.bold = True
        run.font.size = Pt(layout.title_font_size)
        run.font.color.rgb = _rgb(layout.title_color)
        run.font.name = layout.font_family

    def _add_intro_line(self, slide: Slide, spec: SlideSpec, layout: LayoutSpec) -> None:
        """Add the one-line intro sentence below the title (brand header lockup)."""
        if not layout.show_brand_header:
            return
        text = spec.intro_line or spec.key_message
        if not text:
            return
        box = slide.shapes.add_textbox(
            Inches(layout.title_left_inches),
            Inches(layout.intro_top_inches),
            Inches(layout.content_width_inches),
            Inches(0.40),
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(10)
        run.font.color.rgb = _rgb(layout.intro_color)
        run.font.name = layout.font_family

    def _add_subtitle_line(
        self,
        slide: Slide,
        text: str,
        layout: LayoutSpec,
        top_offset: float = 0.0,
    ) -> None:
        box = slide.shapes.add_textbox(
            Inches(layout.content_left_inches),
            Inches(layout.content_top_inches + top_offset),
            Inches(layout.content_width_inches),
            Inches(0.55),
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(layout.subtitle_font_size)
        run.font.italic = True
        run.font.color.rgb = _rgb(layout.accent_color)
        run.font.name = layout.font_family

    def _add_key_message(self, slide: Slide, spec: SlideSpec, layout: LayoutSpec) -> None:
        if not spec.key_message:
            return
        box = slide.shapes.add_textbox(
            Inches(layout.content_left_inches),
            Inches(layout.content_top_inches),
            Inches(layout.content_width_inches),
            Inches(0.50),
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = spec.key_message
        run.font.bold = True
        run.font.size = Pt(layout.subtitle_font_size)
        run.font.color.rgb = _rgb(layout.body_color)
        run.font.name = layout.font_family

    def _add_bullets(self, slide: Slide, spec: SlideSpec, layout: LayoutSpec) -> None:
        """Add bullet points as one native text box with one paragraph per bullet."""
        if not spec.bullets:
            return
        bullet_top = layout.content_top_inches + 0.58
        available_height = layout.content_height_inches - 0.58
        box = slide.shapes.add_textbox(
            Inches(layout.content_left_inches),
            Inches(bullet_top),
            Inches(layout.content_width_inches),
            Inches(max(available_height, 0.5)),
        )
        tf = box.text_frame
        tf.word_wrap = True

        for idx, bullet in enumerate(spec.bullets):
            para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            run = para.add_run()
            run.text = f"\u25aa  {bullet}"
            run.font.size = Pt(layout.body_font_size)
            run.font.color.rgb = _rgb(layout.body_color)
            run.font.name = layout.font_family
            para.space_after = Pt(6)

    def _add_takeaway_bar(self, slide: Slide, spec: SlideSpec, layout: LayoutSpec) -> None:
        """Add the full-width navy bottom bar with the slide's 'so what' sentence."""
        if not spec.takeaway or not layout.show_brand_header:
            return
        # Full-width navy bar
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.0), Inches(layout.takeaway_top_inches),
            Inches(13.33), Inches(layout.takeaway_height_inches),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(layout.takeaway_bg_color)
        bar.line.width = 0

        # Takeaway text centered in the bar
        box = slide.shapes.add_textbox(
            Inches(0.5), Inches(layout.takeaway_top_inches + 0.08),
            Inches(12.33), Inches(layout.takeaway_height_inches - 0.12),
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = spec.takeaway
        run.font.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = _rgb(layout.takeaway_text_color)
        run.font.name = layout.font_family

    def _add_speaker_notes(self, slide: Slide, spec: SlideSpec) -> None:
        """Write speaker_notes into the slide's notes pane."""
        if not spec.speaker_notes:
            return
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = spec.speaker_notes

    def _add_slide_number(self, slide: Slide, number: int, layout: LayoutSpec) -> None:
        if layout.slide_type in (SlideType.TITLE, SlideType.CLOSING):
            return
        box = slide.shapes.add_textbox(
            Inches(12.55), Inches(7.1),
            Inches(0.65), Inches(0.32),
        )
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = str(number)
        run.font.size = Pt(9)
        run.font.color.rgb = _rgb("#AAAAAA")
        run.font.name = layout.font_family

    def _add_footer_label(self, slide: Slide, spec: SlideSpec, layout: LayoutSpec) -> None:
        """Add client/deck label on the left side of the footer row (consulting standard)."""
        if layout.slide_type in (SlideType.TITLE, SlideType.CLOSING, SlideType.SECTION_DIVIDER):
            return
        label = "Rogers Communications  |  AI Governance & Evaluation System"
        box = slide.shapes.add_textbox(
            Inches(0.5), Inches(7.1),
            Inches(9.0), Inches(0.28),
        )
        tf = box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.size = Pt(7.5)
        run.font.color.rgb = _rgb("#AAAAAA")
        run.font.name = layout.font_family

    def _add_stat_boxes(self, slide: Slide, spec: SlideSpec, layout: LayoutSpec) -> None:
        """Render bullets as four KPI stat boxes spanning the full content zone.

        Expected bullet format: "VALUE | LABEL | CONTEXT"
        Example: "85% | Agents Without Governance | No formal validation before deployment"
        """
        bullets = spec.bullets
        if not bullets:
            return

        n = min(len(bullets), 4)
        left_margin = 0.5
        gap = 0.28
        total_w = 13.33 - 2 * left_margin  # 12.33"
        box_w = (total_w - (n - 1) * gap) / n
        box_top = 2.25
        box_h = 3.6

        # Card background tints (alternating warm/cool on the brand peach background)
        fills = ["#FFFFFF", "#FFFFFF", "#FFFFFF", "#FFFFFF"]
        border_color = "#D8D0C8"

        for idx, bullet in enumerate(bullets[:n]):
            parts = [p.strip() for p in str(bullet).split("|")]
            value = parts[0] if parts else "—"
            label = parts[1] if len(parts) > 1 else ""
            context = parts[2] if len(parts) > 2 else ""

            x = left_margin + idx * (box_w + gap)

            # Card background
            card = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(x), Inches(box_top),
                Inches(box_w), Inches(box_h),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = _rgb(fills[idx % len(fills)])
            card.line.color.rgb = _rgb(border_color)
            card.line.width = Pt(0.75)

            # Orange accent top bar
            accent = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(x), Inches(box_top),
                Inches(box_w), Inches(0.07),
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = _rgb(layout.accent_color)
            accent.line.width = 0

            # Large value text
            val_box = slide.shapes.add_textbox(
                Inches(x + 0.18), Inches(box_top + 0.18),
                Inches(box_w - 0.36), Inches(0.9),
            )
            tf = val_box.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = value
            run.font.bold = True
            run.font.size = Pt(36)
            run.font.color.rgb = _rgb(layout.accent_color)
            run.font.name = layout.font_family

            # Label text (bold, navy)
            if label:
                lbl_box = slide.shapes.add_textbox(
                    Inches(x + 0.18), Inches(box_top + 1.15),
                    Inches(box_w - 0.36), Inches(0.55),
                )
                tf = lbl_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = label
                run.font.bold = True
                run.font.size = Pt(11.5)
                run.font.color.rgb = _rgb(layout.title_color)
                run.font.name = layout.font_family

            # Context text (small, grey)
            if context:
                ctx_box = slide.shapes.add_textbox(
                    Inches(x + 0.18), Inches(box_top + 1.78),
                    Inches(box_w - 0.36), Inches(1.5),
                )
                tf = ctx_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = context
                run.font.size = Pt(9.5)
                run.font.color.rgb = _rgb("#5C5C78")
                run.font.name = layout.font_family

        self._add_stat_context_strip(slide, spec, layout)

    # ── Dense consulting helpers ──────────────────────────────────────────────

    def _set_cell_text(
        self,
        cell,
        text: str,
        font_size: float,
        bold: bool,
        color_hex: str,
        font_family: str,
        align=PP_ALIGN.LEFT,
    ) -> None:
        """Set text and font in a PPTX table cell, safely clearing prior runs."""
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        for r_elem in list(p._p.findall(_xml_qn("a:r"))):
            p._p.remove(r_elem)
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.size = Pt(font_size)
        run.font.color.rgb = _rgb(color_hex)
        run.font.name = font_family

    def _parse_bullets_for_dense(self, bullets: list) -> tuple:
        """Split encoded bullets into table data, KPI cards, insight dict, and plain text."""
        table_headers: list = []
        table_rows: list = []
        kpis: list = []
        insight: dict = {}
        regular: list = []
        for bullet in bullets:
            s = str(bullet).strip()
            tag = s[:9].upper()
            if tag.startswith("TABLE:"):
                table_headers = [c.strip() for c in s[6:].split("|") if c.strip()]
            elif tag.startswith("ROW:"):
                table_rows.append([c.strip() for c in s[4:].split("|")])
            elif tag.startswith("KPI:"):
                parts = [p.strip() for p in s[4:].split("|")]
                if parts:
                    kpis.append(parts)
            elif tag.startswith("INSIGHT:"):
                inner = s[8:].strip()
                halves = inner.split("|", 1)
                insight = {
                    "type": halves[0].strip() if halves else "Business Impact",
                    "text": halves[1].strip() if len(halves) > 1 else inner,
                }
            elif s:
                regular.append(s)
        return table_headers, table_rows, kpis, insight, regular

    def _cell_semantic_color(self, val: str) -> Optional[tuple]:
        """Return (bg_hex, text_hex) for semantically meaningful single-word cell values."""
        _MAP = {
            "r":       ("#7B0000", "#FEFBF8"),  # RACI: Responsible
            "a":       ("#E9590C", "#FEFBF8"),  # RACI: Accountable
            "c":       ("#0F3460", "#FEFBF8"),  # RACI: Consulted
            "i":       ("#5C5C78", "#FEFBF8"),  # RACI: Informed
            "high":    ("#7B0000", "#FEFBF8"),
            "medium":  ("#B8520A", "#FEFBF8"),
            "low":     ("#0A5C2A", "#FEFBF8"),
            "pass":    ("#0A5C2A", "#FEFBF8"),
            "fail":    ("#7B0000", "#FEFBF8"),
            "warning": ("#B8520A", "#FEFBF8"),
            "critical":("#5A0000", "#FEFBF8"),
            "at risk": ("#B8520A", "#FEFBF8"),
            "on track":("#0A5C2A", "#FEFBF8"),
        }
        return _MAP.get(val.strip().lower())

    def _add_native_table(
        self,
        slide: Slide,
        headers: list,
        rows: list,
        x: float,
        y: float,
        w: float,
        h: float,
        layout: LayoutSpec,
    ) -> None:
        """PPTX native table: navy header, alternating warm-peach/white data rows."""
        if not headers:
            return
        n_cols = len(headers)
        n_data = max(len(rows), 1)
        tbl_shape = slide.shapes.add_table(
            n_data + 1, n_cols,
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
        tbl = tbl_shape.table
        col_w = Inches(w / n_cols)
        for col in tbl.columns:
            col.width = col_w
        tbl.rows[0].height = Inches(0.33)
        data_h = max(0.25, (h - 0.33) / max(len(rows), 1))
        for i in range(1, n_data + 1):
            tbl.rows[i].height = Inches(data_h)
        # Header: navy bg, off-white bold
        for ci, hdr in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(layout.title_color)
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            self._set_cell_text(cell, hdr, 8.5, True, "#FEFBF8", layout.font_family)
        # Data rows: alternating warm/white, with semantic color override for special values
        fills = ["#F7F0E8", "#FFFFFF"]
        for ri, row_data in enumerate(rows):
            for ci in range(n_cols):
                val = row_data[ci] if ci < len(row_data) else ""
                cell = tbl.cell(ri + 1, ci)
                semantic = self._cell_semantic_color(str(val))
                if semantic:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(semantic[0])
                    cell.margin_left = Inches(0.07)
                    cell.margin_right = Inches(0.05)
                    cell.margin_top = Inches(0.02)
                    cell.margin_bottom = Inches(0.02)
                    self._set_cell_text(cell, str(val), 8.5, True, semantic[1], layout.font_family, PP_ALIGN.CENTER)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(fills[ri % 2])
                    cell.margin_left = Inches(0.07)
                    cell.margin_right = Inches(0.05)
                    cell.margin_top = Inches(0.02)
                    cell.margin_bottom = Inches(0.02)
                    self._set_cell_text(cell, str(val), 8.5, ci == 0, layout.body_color, layout.font_family)

    def _add_executive_callout(
        self,
        slide: Slide,
        insight: dict,
        x: float,
        y: float,
        w: float,
        h: float,
        layout: LayoutSpec,
    ) -> float:
        """Dark callout box with accent-colored label. Returns y-coordinate below the box."""
        panel_type = insight.get("type", "Business Impact")
        text = insight.get("text", "")
        bg_color, accent = "#0C1445", layout.accent_color
        p_lower = panel_type.lower()
        if "risk" in p_lower:
            bg_color, accent = "#3D0D0D", "#FF6B55"
        elif "success" in p_lower or "criteria" in p_lower:
            bg_color, accent = "#0A3320", "#34C97A"
        bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb(bg_color)
        bg.line.width = 0
        strip = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.07),
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = _rgb(accent)
        strip.line.width = 0
        lbl = slide.shapes.add_textbox(
            Inches(x + 0.12), Inches(y + 0.10), Inches(w - 0.24), Inches(0.28),
        )
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = panel_type.upper()
        run.font.bold = True
        run.font.size = Pt(7.5)
        run.font.color.rgb = _rgb(accent)
        run.font.name = layout.font_family
        if text:
            content = slide.shapes.add_textbox(
                Inches(x + 0.12), Inches(y + 0.41), Inches(w - 0.24), Inches(h - 0.47),
            )
            tf = content.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(9)
            run.font.color.rgb = _rgb("#D8D8E4")
            run.font.name = layout.font_family
        return y + h + 0.13

    def _add_kpi_mini_card(
        self,
        slide: Slide,
        kpi_parts: list,
        x: float,
        y: float,
        w: float,
        layout: LayoutSpec,
    ) -> float:
        """Compact KPI card: large value left, label/context right. Returns next y."""
        h = 0.70
        value   = kpi_parts[0] if kpi_parts else "—"
        label   = kpi_parts[1] if len(kpi_parts) > 1 else ""
        context = kpi_parts[2] if len(kpi_parts) > 2 else ""
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb("#FFFFFF")
        card.line.color.rgb = _rgb("#D5CCC4")
        card.line.width = Pt(0.5)
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(layout.accent_color)
        bar.line.width = 0
        val_box = slide.shapes.add_textbox(
            Inches(x + 0.12), Inches(y + 0.04), Inches(w * 0.44), Inches(h - 0.08),
        )
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = value
        run.font.bold = True
        run.font.size = Pt(20)
        run.font.color.rgb = _rgb(layout.accent_color)
        run.font.name = layout.font_family
        lc_box = slide.shapes.add_textbox(
            Inches(x + w * 0.46), Inches(y + 0.04), Inches(w * 0.54 - 0.06), Inches(h - 0.08),
        )
        tf = lc_box.text_frame
        tf.word_wrap = True
        if label:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = label
            run.font.bold = True
            run.font.size = Pt(8.5)
            run.font.color.rgb = _rgb(layout.title_color)
            run.font.name = layout.font_family
        if context:
            p2 = tf.add_paragraph() if label else tf.paragraphs[0]
            run2 = p2.add_run()
            run2.text = context
            run2.font.size = Pt(7.5)
            run2.font.color.rgb = _rgb("#6C6C88")
            run2.font.name = layout.font_family
        return y + h + 0.10

    def _add_vertical_divider(
        self, slide: Slide, x: float, top: float, bottom: float, layout: LayoutSpec,
    ) -> None:
        """Thin accent-colored vertical separator between content zones."""
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(top), Inches(0.02), Inches(bottom - top),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = _rgb(layout.accent_color)
        line.line.width = 0

    def _fill_left_gap(self, slide: Slide, spec: SlideSpec, start_y: float, layout: LayoutSpec) -> None:
        """Fill unused left-zone space with supporting context drawn from intro_line / takeaway."""
        avail = _DC_BOTTOM - start_y - 0.05
        if avail < 0.45:
            return
        sep = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(_DC_LEFT_X), Inches(start_y), Inches(_DC_LEFT_W), Inches(0.02),
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = _rgb("#D5CCC4")
        sep.line.width = 0
        start_y += 0.07
        avail  -= 0.07
        lbl = slide.shapes.add_textbox(
            Inches(_DC_LEFT_X), Inches(start_y), Inches(_DC_LEFT_W), Inches(0.22),
        )
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "SUPPORTING CONTEXT"
        run.font.bold = True
        run.font.size = Pt(7.5)
        run.font.color.rgb = _rgb(layout.accent_color)
        run.font.name = layout.font_family
        start_y += 0.25
        avail  -= 0.25
        if avail < 0.18:
            return
        lines = []
        if spec.intro_line:
            lines.append((f"\u25b8  {spec.intro_line}", False))
        if spec.takeaway and spec.takeaway != spec.key_message:
            lines.append((f"\u25b8  {spec.takeaway}", True))
        if lines:
            box = slide.shapes.add_textbox(
                Inches(_DC_LEFT_X), Inches(start_y), Inches(_DC_LEFT_W), Inches(avail),
            )
            tf = box.text_frame
            tf.word_wrap = True
            for i, (text, bold) in enumerate(lines[:3]):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = para.add_run()
                run.text = text
                run.font.size = Pt(9)
                run.font.bold = bold
                run.font.color.rgb = _rgb(layout.body_color)
                run.font.name = layout.font_family
                para.space_after = Pt(5)

    def _fill_right_gap(self, slide: Slide, spec: SlideSpec, start_y: float, layout: LayoutSpec) -> None:
        """Fill unused right-panel space with a warm section-objective box."""
        avail = _DC_BOTTOM - start_y - 0.05
        if avail < 0.45:
            return
        bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(_DC_RIGHT_X), Inches(start_y), Inches(_DC_RIGHT_W), Inches(avail),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb("#F5EDE5")
        bg.line.color.rgb = _rgb("#D5CCC4")
        bg.line.width = Pt(0.5)
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(_DC_RIGHT_X), Inches(start_y), Inches(0.05), Inches(avail),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(layout.accent_color)
        bar.line.width = 0
        lbl = slide.shapes.add_textbox(
            Inches(_DC_RIGHT_X + 0.12), Inches(start_y + 0.08),
            Inches(_DC_RIGHT_W - 0.17), Inches(0.22),
        )
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "SECTION OBJECTIVE"
        run.font.bold = True
        run.font.size = Pt(7.5)
        run.font.color.rgb = _rgb(layout.accent_color)
        run.font.name = layout.font_family
        text = spec.intro_line or spec.key_message or ""
        if text and avail > 0.5:
            content = slide.shapes.add_textbox(
                Inches(_DC_RIGHT_X + 0.12), Inches(start_y + 0.32),
                Inches(_DC_RIGHT_W - 0.17), Inches(avail - 0.38),
            )
            tf = content.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(8.5)
            run.font.color.rgb = _rgb(layout.body_color)
            run.font.name = layout.font_family

    def _add_stat_context_strip(self, slide: Slide, spec: SlideSpec, layout: LayoutSpec) -> None:
        """Narrow context strip below stat boxes, filling the gap before the takeaway bar."""
        text = spec.intro_line or spec.key_message or ""
        if not text:
            return
        strip_y, strip_h = 5.90, 0.42
        bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.5), Inches(strip_y), Inches(12.33), Inches(strip_h),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb("#F0E8DF")
        bg.line.color.rgb = _rgb(layout.accent_color)
        bg.line.width = Pt(0.5)
        lbl = slide.shapes.add_textbox(
            Inches(0.62), Inches(strip_y + 0.07), Inches(1.4), Inches(strip_h - 0.14),
        )
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "CONTEXT"
        run.font.bold = True
        run.font.size = Pt(7.5)
        run.font.color.rgb = _rgb(layout.accent_color)
        run.font.name = layout.font_family
        sep = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(2.12), Inches(strip_y + 0.08), Inches(0.02), Inches(strip_h - 0.16),
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = _rgb(layout.accent_color)
        sep.line.width = 0
        txt = slide.shapes.add_textbox(
            Inches(2.26), Inches(strip_y + 0.07), Inches(10.45), Inches(strip_h - 0.14),
        )
        tf = txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(9)
        run.font.color.rgb = _rgb(layout.body_color)
        run.font.name = layout.font_family

    def _detect_dense_pattern(self, bullets: list) -> str:
        """Return which visual pattern the bullet encoding requests."""
        for b in bullets:
            s = str(b).strip().upper()
            if s.startswith("COL1:") or s.startswith("COL1_ITEM:"):
                return "three_column"
            if s.startswith("COMPARE_LEFT:") or s.startswith("BEFORE:"):
                return "two_column"
            if s.startswith("STEP:"):
                return "process_steps"
        return "standard"

    def _render_two_column(self, slide: Slide, spec: SlideSpec, layout: LayoutSpec) -> None:
        """Before/After or Current/Target two-column comparison layout."""
        left_hdr, right_hdr = "CURRENT STATE", "TARGET STATE"
        rows: list = []
        kpis: list = []
        insight: dict = {}
        regular: list = []
        for b in spec.bullets:
            s = str(b).strip(); u = s.upper()
            if u.startswith("COMPARE_LEFT:") or u.startswith("BEFORE:"):
                left_hdr = s.split(":", 1)[1].strip().upper()
            elif u.startswith("COMPARE_RIGHT:") or u.startswith("AFTER:"):
                right_hdr = s.split(":", 1)[1].strip().upper()
            elif u.startswith("COMPARE_ROW:"):
                parts = s[12:].split("|", 1)
                rows.append((parts[0].strip(), parts[1].strip() if len(parts) > 1 else ""))
            elif u.startswith("KPI:"):
                kpis.append([p.strip() for p in s[4:].split("|")])
            elif u.startswith("INSIGHT:"):
                inner = s[8:].strip(); halves = inner.split("|", 1)
                insight = {"type": halves[0].strip(), "text": halves[1].strip() if len(halves) > 1 else inner}
            elif s and not u.startswith("TABLE:") and not u.startswith("ROW:"):
                regular.append(s)
        if not rows and regular:
            mid = (len(regular) + 1) // 2
            rows = [(regular[i], regular[i + mid] if i + mid < len(regular) else "") for i in range(mid)]

        lx = 0.50; full_w = 12.33
        col_w = (full_w - 0.50) / 2  # 5.915"
        rx = lx + col_w + 0.50
        y = _DC_TOP

        if spec.key_message:
            km = slide.shapes.add_textbox(Inches(lx), Inches(y), Inches(full_w), Inches(0.35))
            tf = km.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; run = p.add_run()
            run.text = f"\u258c  {spec.key_message}"
            run.font.bold = True; run.font.size = Pt(9.5)
            run.font.color.rgb = _rgb(layout.title_color); run.font.name = layout.font_family
            y += 0.42

        insight_h = 0.68 if (insight or spec.takeaway) else 0.0
        col_h = _DC_BOTTOM - y - insight_h - 0.10

        for ci, (cx, hdr_bg, hdr_acc, hdr_text) in enumerate([
            (lx,  "#0F0F37", layout.accent_color, left_hdr),
            (rx,  "#0C3320", "#34C97A",            right_hdr),
        ]):
            hdr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(cx), Inches(y), Inches(col_w), Inches(0.38))
            hdr.fill.solid(); hdr.fill.fore_color.rgb = _rgb(hdr_bg); hdr.line.width = 0
            ht = slide.shapes.add_textbox(
                Inches(cx + 0.12), Inches(y + 0.06), Inches(col_w - 0.24), Inches(0.28))
            tf = ht.text_frame; p = tf.paragraphs[0]; run = p.add_run()
            run.text = hdr_text; run.font.bold = True; run.font.size = Pt(9)
            run.font.color.rgb = _rgb(hdr_acc); run.font.name = layout.font_family

            row_y = y + 0.42
            # expand rows to fill the full column height — no empty space below last row
            row_h = (col_h - 0.42) / max(len(rows), 1)
            for ri, pair in enumerate(rows[:12]):
                val = pair[ci] if ci < len(pair) else ""
                rb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    Inches(cx), Inches(row_y), Inches(col_w), Inches(row_h))
                rb.fill.solid()
                rb.fill.fore_color.rgb = _rgb("#F7F0E8" if ri % 2 == 0 else "#FFFFFF")
                rb.line.width = 0
                if val:
                    tb = slide.shapes.add_textbox(
                        Inches(cx + 0.10), Inches(row_y + 0.04), Inches(col_w - 0.20), Inches(row_h - 0.08))
                    tf = tb.text_frame; tf.word_wrap = True
                    p = tf.paragraphs[0]; run = p.add_run()
                    run.text = val; run.font.size = Pt(9)
                    run.font.color.rgb = _rgb(layout.body_color); run.font.name = layout.font_family
                row_y += row_h

        if insight_h > 0:
            eff = insight if insight else {"type": "Key Takeaway", "text": spec.takeaway or spec.key_message or ""}
            strip_y = y + col_h + 0.05
            avail = _DC_BOTTOM - strip_y - 0.02
            if eff.get("text") and avail > 0.30:
                self._add_executive_callout(slide, eff, lx, strip_y, full_w, avail, layout)

    def _render_three_column(self, slide: Slide, spec: SlideSpec, layout: LayoutSpec) -> None:
        """Three-pillar framework: three equal columns with colored headers and bullet content."""
        cols = [{"header": f"PILLAR {i+1}", "items": []} for i in range(3)]
        insight: dict = {}; regular: list = []
        for b in spec.bullets:
            s = str(b).strip(); u = s[:14].upper()
            if   u.startswith("COL1:"):      cols[0]["header"] = s[5:].strip().upper()
            elif u.startswith("COL2:"):      cols[1]["header"] = s[5:].strip().upper()
            elif u.startswith("COL3:"):      cols[2]["header"] = s[5:].strip().upper()
            elif u.startswith("COL1_ITEM:"): cols[0]["items"].append(s[10:].strip())
            elif u.startswith("COL2_ITEM:"): cols[1]["items"].append(s[10:].strip())
            elif u.startswith("COL3_ITEM:"): cols[2]["items"].append(s[10:].strip())
            elif u.startswith("INSIGHT:"):
                inner = s[8:].strip(); halves = inner.split("|", 1)
                insight = {"type": halves[0].strip(), "text": halves[1].strip() if len(halves) > 1 else inner}
            elif s and not u.startswith("KPI:") and not u.startswith("TABLE:") and not u.startswith("ROW:"):
                regular.append(s)
        if all(not c["items"] for c in cols) and regular:
            per = max(1, (len(regular) + 2) // 3)
            for i, item in enumerate(regular):
                cols[min(i // per, 2)]["items"].append(item)

        lx = 0.50; full_w = 12.33; gap = 0.35; n = 3
        col_w = (full_w - (n - 1) * gap) / n
        y = _DC_TOP

        if spec.key_message:
            km = slide.shapes.add_textbox(Inches(lx), Inches(y), Inches(full_w), Inches(0.35))
            tf = km.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; run = p.add_run()
            run.text = f"\u258c  {spec.key_message}"
            run.font.bold = True; run.font.size = Pt(9.5)
            run.font.color.rgb = _rgb(layout.title_color); run.font.name = layout.font_family
            y += 0.42

        insight_h = 0.68 if insight else 0.0
        col_area_h = _DC_BOTTOM - y - insight_h - 0.10
        hdr_bgs  = [layout.title_color, "#1A4A7A", "#0A3320"]
        hdr_accs = [layout.accent_color, "#4DA6FF", "#34C97A"]

        for ci, col in enumerate(cols):
            cx = lx + ci * (col_w + gap)
            hb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(cx), Inches(y), Inches(col_w), Inches(0.40))
            hb.fill.solid(); hb.fill.fore_color.rgb = _rgb(hdr_bgs[ci]); hb.line.width = 0
            st = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(cx), Inches(y), Inches(col_w), Inches(0.06))
            st.fill.solid(); st.fill.fore_color.rgb = _rgb(hdr_accs[ci]); st.line.width = 0
            ht = slide.shapes.add_textbox(
                Inches(cx + 0.10), Inches(y + 0.08), Inches(col_w - 0.20), Inches(0.32))
            tf = ht.text_frame; p = tf.paragraphs[0]; run = p.add_run()
            run.text = col["header"]; run.font.bold = True; run.font.size = Pt(9)
            run.font.color.rgb = _rgb("#FEFBF8"); run.font.name = layout.font_family
            body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(cx), Inches(y + 0.40), Inches(col_w), Inches(col_area_h - 0.40))
            body.fill.solid(); body.fill.fore_color.rgb = _rgb("#FAFAF8")
            body.line.color.rgb = _rgb("#D5CCC4"); body.line.width = Pt(0.5)
            # Items: distribute vertically with even spacing; reserve bottom for footer
            n_items = len(col["items"][:8])
            body_avail = col_area_h - 0.48
            footer_budget = 0.52 if body_avail > max(n_items * 0.30 + 0.55, 1.0) else 0.0
            items_budget = body_avail - footer_budget
            extra_h = max(0.0, items_budget - n_items * 0.22)
            even_spacing = Pt(max(5, min(30, extra_h * 72 / max(n_items, 1))))
            ib = slide.shapes.add_textbox(
                Inches(cx + 0.12), Inches(y + 0.48), Inches(col_w - 0.24), Inches(items_budget))
            tf = ib.text_frame; tf.word_wrap = True
            for ii, item in enumerate(col["items"][:8]):
                para = tf.paragraphs[0] if ii == 0 else tf.add_paragraph()
                run = para.add_run()
                run.text = f"\u25ba  {item}"
                run.font.size = Pt(9); run.font.color.rgb = _rgb(layout.body_color)
                run.font.name = layout.font_family; para.space_after = even_spacing
            # Column footer: supplementary context to fill remaining space
            if footer_budget >= 0.45:
                fy = y + 0.48 + items_budget
                fsep = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    Inches(cx + 0.08), Inches(fy), Inches(col_w - 0.16), Inches(0.02))
                fsep.fill.solid(); fsep.fill.fore_color.rgb = _rgb(hdr_accs[ci]); fsep.line.width = 0
                flbl = slide.shapes.add_textbox(
                    Inches(cx + 0.10), Inches(fy + 0.05), Inches(col_w - 0.20), Inches(0.20))
                tf = flbl.text_frame; p = tf.paragraphs[0]; run = p.add_run()
                run.text = ["KEY INSIGHT", "SUPPORTING CONTEXT", "KEY IMPLICATION"][ci % 3]
                run.font.bold = True; run.font.size = Pt(7.5)
                run.font.color.rgb = _rgb(hdr_accs[ci]); run.font.name = layout.font_family
                ftext = ([spec.takeaway, spec.intro_line, spec.key_message][ci % 3] or "")[:130]
                if ftext and footer_budget > 0.48:
                    ftxt = slide.shapes.add_textbox(
                        Inches(cx + 0.10), Inches(fy + 0.27),
                        Inches(col_w - 0.20), Inches(footer_budget - 0.30))
                    tf = ftxt.text_frame; tf.word_wrap = True
                    p = tf.paragraphs[0]; run = p.add_run(); run.text = ftext
                    run.font.size = Pt(8); run.font.color.rgb = _rgb(layout.body_color)
                    run.font.name = layout.font_family

        if insight:
            sy = y + col_area_h + 0.05
            avail = _DC_BOTTOM - sy - 0.02
            if avail > 0.30:
                self._add_executive_callout(slide, insight, lx, sy, full_w, avail, layout)

    def _render_process_steps(self, slide: Slide, spec: SlideSpec, layout: LayoutSpec) -> None:
        """Horizontal numbered process-step boxes with arrows between them."""
        steps: list = []; kpis: list = []; insight: dict = {}; regular: list = []
        for b in spec.bullets:
            s = str(b).strip(); u = s[:10].upper()
            if u.startswith("STEP:"):
                parts = s[5:].split("|")
                steps.append({"num": parts[0].strip(), "name": parts[1].strip() if len(parts)>1 else "",
                               "desc": parts[2].strip() if len(parts)>2 else ""})
            elif u.startswith("KPI:"):
                kpis.append([p.strip() for p in s[4:].split("|")])
            elif u.startswith("INSIGHT:"):
                inner = s[8:].strip(); halves = inner.split("|", 1)
                insight = {"type": halves[0].strip(), "text": halves[1].strip() if len(halves)>1 else inner}
            elif s:
                regular.append(s)
        if not steps and regular:
            steps = [{"num": str(i+1), "name": f"Step {i+1}", "desc": r} for i, r in enumerate(regular[:5])]

        lx = 0.50; full_w = 12.33; n = max(len(steps), 1); gap = 0.30
        box_w = (full_w - (n - 1) * gap) / n
        y = _DC_TOP

        if spec.key_message:
            km = slide.shapes.add_textbox(Inches(lx), Inches(y), Inches(full_w), Inches(0.35))
            tf = km.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; run = p.add_run()
            run.text = f"\u258c  {spec.key_message}"
            run.font.bold = True; run.font.size = Pt(9.5)
            run.font.color.rgb = _rgb(layout.title_color); run.font.name = layout.font_family
            y += 0.42

        bottom_h = 0.68 if (insight or kpis) else 0.0
        box_h = _DC_BOTTOM - y - bottom_h - 0.12

        for si, step in enumerate(steps[:5]):
            bx = lx + si * (box_w + gap)
            bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(bx), Inches(y), Inches(box_w), Inches(box_h))
            bg.fill.solid(); bg.fill.fore_color.rgb = _rgb("#FAFAF8")
            bg.line.color.rgb = _rgb("#D5CCC4"); bg.line.width = Pt(0.75)
            ac = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(bx), Inches(y), Inches(box_w), Inches(0.06))
            ac.fill.solid(); ac.fill.fore_color.rgb = _rgb(layout.accent_color); ac.line.width = 0
            nb = slide.shapes.add_textbox(
                Inches(bx + 0.10), Inches(y + 0.08), Inches(box_w - 0.20), Inches(0.55))
            tf = nb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = step["num"]
            run.font.bold = True; run.font.size = Pt(22)
            run.font.color.rgb = _rgb(layout.accent_color); run.font.name = layout.font_family
            if step["name"]:
                nmb = slide.shapes.add_textbox(
                    Inches(bx + 0.10), Inches(y + 0.66), Inches(box_w - 0.20), Inches(0.40))
                tf = nmb.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                run = p.add_run(); run.text = step["name"]
                run.font.bold = True; run.font.size = Pt(9.5)
                run.font.color.rgb = _rgb(layout.title_color); run.font.name = layout.font_family
            if step["desc"]:
                # Reserve 0.44" at bottom for gate footer (only when box is tall enough)
                gate_budget = 0.44 if box_h > 1.72 else 0.0
                db = slide.shapes.add_textbox(
                    Inches(bx + 0.10), Inches(y + 1.12), Inches(box_w - 0.20), Inches(box_h - 1.22 - gate_budget))
                tf = db.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; run = p.add_run(); run.text = step["desc"]
                run.font.size = Pt(9); run.font.color.rgb = _rgb("#5C5C78"); run.font.name = layout.font_family
            else:
                gate_budget = 0.44 if box_h > 1.72 else 0.0
            # Gate criteria footer inside every step box
            if box_h > 1.72:
                gy = y + box_h - 0.42
                gsep = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    Inches(bx + 0.08), Inches(gy), Inches(box_w - 0.16), Inches(0.02))
                gsep.fill.solid(); gsep.fill.fore_color.rgb = _rgb("#D5CCC4"); gsep.line.width = 0
                glbl = slide.shapes.add_textbox(
                    Inches(bx + 0.10), Inches(gy + 0.04), Inches(box_w - 0.20), Inches(0.16))
                tf = glbl.text_frame; p = tf.paragraphs[0]; run = p.add_run()
                run.text = "GATE CRITERIA"; run.font.bold = True; run.font.size = Pt(6.5)
                run.font.color.rgb = _rgb(layout.accent_color); run.font.name = layout.font_family
                criteria = ((step["desc"].split(".")[0] + ".") if step["desc"] else step["name"])[:65]
                if criteria:
                    gval = slide.shapes.add_textbox(
                        Inches(bx + 0.10), Inches(gy + 0.22), Inches(box_w - 0.20), Inches(0.17))
                    tf = gval.text_frame; p = tf.paragraphs[0]; run = p.add_run()
                    run.text = criteria; run.font.size = Pt(7.5)
                    run.font.color.rgb = _rgb("#5C5C78"); run.font.name = layout.font_family
            if si < len(steps) - 1:
                ab = slide.shapes.add_textbox(
                    Inches(bx + box_w), Inches(y + box_h / 2 - 0.15), Inches(gap), Inches(0.30))
                tf = ab.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                run = p.add_run(); run.text = "\u2192"
                run.font.bold = True; run.font.size = Pt(14)
                run.font.color.rgb = _rgb(layout.accent_color); run.font.name = layout.font_family

        if bottom_h > 0:
            sy = y + box_h + 0.06; avail = _DC_BOTTOM - sy - 0.02
            eff = insight if insight else {"type": "Key Takeaway", "text": spec.takeaway or spec.key_message or ""}
            if kpis and eff.get("text"):
                ins_w = full_w * 0.54
                self._add_executive_callout(slide, eff, lx, sy, ins_w, avail, layout)
                kx = lx + ins_w + 0.10; kw = full_w - ins_w - 0.10
                card_w = (kw - (len(kpis[:3]) - 1) * 0.08) / max(len(kpis[:3]), 1)
                for ki, kp in enumerate(kpis[:3]):
                    self._add_kpi_mini_card(slide, kp, kx + ki * (card_w + 0.08), sy, card_w, layout)
            elif eff.get("text") and avail > 0.30:
                self._add_executive_callout(slide, eff, lx, sy, full_w, avail, layout)

    def _render_standard_dense(self, slide: Slide, spec: SlideSpec, layout: LayoutSpec) -> None:
        """Standard dense consulting: left analysis zone + right executive panel + KPI cards."""
        table_headers, table_rows, kpis, insight, regular = self._parse_bullets_for_dense(spec.bullets)

        self._add_vertical_divider(slide, _DC_DIV_X, _DC_TOP, _DC_BOTTOM, layout)

        left_y = _DC_TOP
        if spec.key_message:
            km_box = slide.shapes.add_textbox(
                Inches(_DC_LEFT_X), Inches(left_y), Inches(_DC_LEFT_W), Inches(0.36))
            tf = km_box.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; run = p.add_run()
            run.text = f"\u258c  {spec.key_message}"
            run.font.bold = True; run.font.size = Pt(9.5)
            run.font.color.rgb = _rgb(layout.title_color); run.font.name = layout.font_family
            left_y += 0.40

        if table_headers and table_rows:
            lbl = slide.shapes.add_textbox(
                Inches(_DC_LEFT_X), Inches(left_y), Inches(_DC_LEFT_W), Inches(0.22))
            tf = lbl.text_frame; p = tf.paragraphs[0]; run = p.add_run()
            run.text = "ANALYSIS"; run.font.bold = True; run.font.size = Pt(7.5)
            run.font.color.rgb = _rgb(layout.accent_color); run.font.name = layout.font_family
            left_y += 0.23
            available = _DC_BOTTOM - left_y - 0.05
            n_rows = len(table_rows)
            natural_h  = (n_rows + 1) * 0.31
            expanded_h = min(available * 0.78, (n_rows + 1) * 0.52)
            table_h = min(available - 0.05, max(natural_h, expanded_h))
            self._add_native_table(slide, table_headers, table_rows,
                                   _DC_LEFT_X, left_y, _DC_LEFT_W, table_h, layout)
            left_y += table_h + 0.10

        if regular and left_y < _DC_BOTTOM - 0.4:
            est_h = min(len(regular[:6]) * 0.28, _DC_BOTTOM - left_y - 0.1)
            bul_box = slide.shapes.add_textbox(
                Inches(_DC_LEFT_X), Inches(left_y), Inches(_DC_LEFT_W), Inches(est_h))
            tf = bul_box.text_frame; tf.word_wrap = True
            for i, bullet in enumerate(regular[:6]):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = para.add_run()
                run.text = f"\u25b8  {bullet}"; run.font.size = Pt(9.5)
                run.font.color.rgb = _rgb(layout.body_color); run.font.name = layout.font_family
                para.space_after = Pt(4)
            left_y += est_h + 0.08

        self._fill_left_gap(slide, spec, left_y, layout)

        right_y = _DC_TOP
        effective_insight = insight if insight else {"type": "Executive Insight", "text": spec.key_message or ""}
        if effective_insight.get("text"):
            callout_h = min(1.78, (_DC_BOTTOM - _DC_TOP) * 0.43)
            right_y = self._add_executive_callout(
                slide, effective_insight, _DC_RIGHT_X, right_y, _DC_RIGHT_W, callout_h, layout)
        for kpi_parts in kpis[:3]:
            if right_y + 0.70 <= _DC_BOTTOM + 0.05:
                right_y = self._add_kpi_mini_card(
                    slide, kpi_parts, _DC_RIGHT_X, right_y, _DC_RIGHT_W, layout)
        self._fill_right_gap(slide, spec, right_y, layout)

    def _add_dense_content(self, slide: Slide, spec: SlideSpec, layout: LayoutSpec) -> None:
        """Route to the correct dense pattern based on bullet-prefix encoding."""
        pattern = self._detect_dense_pattern(spec.bullets)
        if pattern == "two_column":
            self._render_two_column(slide, spec, layout)
        elif pattern == "three_column":
            self._render_three_column(slide, spec, layout)
        elif pattern == "process_steps":
            self._render_process_steps(slide, spec, layout)
        else:
            self._render_standard_dense(slide, spec, layout)

    def _add_insight_strip(self, slide: Slide, spec: SlideSpec, layout: LayoutSpec) -> None:
        """Narrow 'Key Insight' bar rendered below the diagram on visual_dominant slides."""
        strip_y = 5.78
        strip_h = 0.52
        bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.5), Inches(strip_y), Inches(12.33), Inches(strip_h),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb("#EEE6DC")
        bg.line.color.rgb = _rgb(layout.accent_color)
        bg.line.width = Pt(0.5)
        lbl = slide.shapes.add_textbox(
            Inches(0.62), Inches(strip_y + 0.07), Inches(1.25), Inches(strip_h - 0.14),
        )
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "KEY INSIGHT"
        run.font.bold = True
        run.font.size = Pt(7.5)
        run.font.color.rgb = _rgb(layout.accent_color)
        run.font.name = layout.font_family
        sep = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(1.98), Inches(strip_y + 0.08), Inches(0.02), Inches(strip_h - 0.16),
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = _rgb(layout.accent_color)
        sep.line.width = 0
        text = spec.key_message or spec.takeaway or ""
        if text:
            txt = slide.shapes.add_textbox(
                Inches(2.12), Inches(strip_y + 0.07), Inches(10.55), Inches(strip_h - 0.14),
            )
            tf = txt.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(9.5)
            run.font.color.rgb = _rgb(layout.body_color)
            run.font.name = layout.font_family

    def render(
        self,
        slide: Slide,
        spec: SlideSpec,
        layout: LayoutSpec,
    ) -> None:
        """Apply all text layers to *slide* following the ML arteka brand header lockup."""
        self._set_background(slide, layout)
        self._add_eyebrow(slide, spec, layout)
        self._add_tick_rule(slide, layout)
        self._add_title(slide, spec, layout)
        self._add_intro_line(slide, spec, layout)

        variant = getattr(spec, "layout_variant", "split")

        if spec.slide_type in (SlideType.TITLE, SlideType.CLOSING):
            if spec.subtitle:
                self._add_subtitle_line(slide, spec.subtitle, layout)
            elif spec.key_message:
                self._add_subtitle_line(slide, spec.key_message, layout)
        elif spec.slide_type in (SlideType.AGENDA, SlideType.SECTION_DIVIDER):
            self._add_key_message(slide, spec, layout)
            self._add_bullets(slide, spec, layout)
        elif variant == "stat_band":
            self._add_stat_boxes(slide, spec, layout)
            self._add_takeaway_bar(slide, spec, layout)
        elif variant == "dense_consulting":
            self._add_dense_content(slide, spec, layout)
            self._add_takeaway_bar(slide, spec, layout)
        else:
            # visual_dominant, split, content_heavy
            if layout.visual_dominant:
                self._add_insight_strip(slide, spec, layout)
            else:
                self._add_key_message(slide, spec, layout)
                self._add_bullets(slide, spec, layout)
            self._add_takeaway_bar(slide, spec, layout)

        self._add_footer_label(slide, spec, layout)
        self._add_slide_number(slide, spec.slide_number, layout)
        self._add_speaker_notes(slide, spec)
