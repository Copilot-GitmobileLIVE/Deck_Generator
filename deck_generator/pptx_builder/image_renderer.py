"""
image_renderer.py — ImageRenderer

Responsibility:
    Insert one image file (PNG) into a python-pptx Slide at the exact position
    and size defined by a LayoutSpec.

This class is intentionally simple — it has a single public method and
no state.  All positioning decisions live in LayoutSpec (produced by
LayoutAgent), keeping the renderer free of layout logic.

IMPORTANT: The z-order of the image in the final PPTX depends on when this
method is called relative to other shape additions.  PPTBuilder controls
the call order explicitly (see builder.py for details).
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

from pptx.slide import Slide
from pptx.util import Inches

from deck_generator.models import LayoutSpec

logger = logging.getLogger("deck_generator.image_renderer")


class ImageRenderer:
    """Adds a generated image to a slide at the position defined by LayoutSpec."""

    def render(
        self,
        slide: Slide,
        layout: LayoutSpec,
        image_path: Optional[str],
    ) -> None:
        """Insert *image_path* into *slide* at the coordinates from *layout*.

        Args:
            slide:      The python-pptx Slide object to add the image to.
            layout:     Provides image_left, image_top, image_width, image_height
                        (all in inches).
            image_path: Filesystem path to the PNG file.  If None or the file
                        does not exist on disk, the method silently returns so
                        the rest of the slide still renders correctly.
        """
        if not image_path:
            return  # No image was selected for this slide — skip silently
        if not Path(image_path).exists():
            logger.warning("ImageRenderer: file not found — %s", image_path)
            return
        if layout.image_width_inches <= 0 or layout.image_height_inches <= 0:
            return  # stat_band and other image-free layouts

        try:
            # add_picture(image_file, left, top, width, height)
            # All positional arguments must be in EMU; Inches() converts for us.
            slide.shapes.add_picture(
                image_path,
                Inches(layout.image_left_inches),
                Inches(layout.image_top_inches),
                Inches(layout.image_width_inches),
                Inches(layout.image_height_inches),
            )
        except Exception as exc:
            # Log and continue — a broken image should not crash the whole deck.
            logger.error("ImageRenderer: failed to add picture — %s", exc)
