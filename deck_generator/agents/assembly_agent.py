"""
assembly_agent.py — Assembly Agent

Responsibility:
    Take all pipeline outputs that are now in DeckState (slides, layout specs,
    selected images) and produce the final .pptx file on disk.

This agent is a thin orchestration wrapper around PPTBuilder.  It is
responsible for:
    1. Deriving a descriptive output filename from the brief title + timestamp.
    2. Converting list-based state fields into the dict-keyed structures that
       PPTBuilder expects.
    3. Calling PPTBuilder.build() and storing the resulting path in state.

No LLM calls are made here — this is pure file I/O.
"""
from __future__ import annotations

import logging
from datetime import datetime
from pathlib import Path
from typing import Dict

from pptx import Presentation
from pptx.util import Inches

from deck_generator.config import get_settings
from deck_generator.models import DeckState, LayoutSpec, SelectedImage, SlideType
from deck_generator.pptx_builder.builder import PPTBuilder

# ML arteka logo specs per slide type: (logo_filename, x, y, w, h)
_ML_LOGOS: Dict[str, tuple] = {
    "light": ("mlarteka-logo.png",      11.78, 0.30, 1.05, 0.23),
    "title": ("mlarteka-logo-dark.png", 10.96, 0.30, 1.85, 0.41),
    "divider": ("mlarteka-logo-dark.png", 11.78, 0.30, 1.05, 0.23),
    "closing": ("mlarteka-logo-dark.png",  5.74, 1.60, 1.85, 0.41),
}

_ASSETS_DIR = Path(__file__).parent.parent / "assets"

logger = logging.getLogger("deck_generator.assembly_agent")


class AssemblyAgent:
    """Drives the PPTBuilder with slides, layouts, and selected images.

    Filename convention:
        <sanitised_title>_<YYYYMMDD_HHMMSS>.pptx
    The timestamp prevents accidental overwrites when the pipeline is run
    multiple times for the same brief.
    """

    def _add_brand_logos(
        self,
        pptx_path: str,
        slides: list,
        layouts: Dict[int, LayoutSpec],
    ) -> None:
        """Overlay the ML arteka logo on every slide at brand-specified coordinates.

        Logo variant and position differ by slide type per the brand spec:
          - Light slides (CONTENT, AGENDA): full-colour logo, top-right
          - Title/Section divider (dark): dark-variant logo, top-right
          - Closing (dark, centered): dark-variant logo, centred below title
        Skips gracefully if the assets directory or a logo file is missing.
        """
        if not _ASSETS_DIR.exists():
            logger.warning("AssemblyAgent: brand assets directory not found — skipping logo placement")
            return

        _type_to_key = {
            SlideType.TITLE: "title",
            SlideType.SECTION_DIVIDER: "divider",
            SlideType.CLOSING: "closing",
        }

        prs = Presentation(pptx_path)
        slide_map = {s.slide_number: s for s in slides}

        for idx, pptx_slide in enumerate(prs.slides):
            slide_number = idx + 1
            spec = slide_map.get(slide_number)
            if spec is None:
                continue

            logo_key = _type_to_key.get(spec.slide_type, "light")
            filename, lx, ly, lw, lh = _ML_LOGOS[logo_key]
            logo_path = _ASSETS_DIR / filename

            if not logo_path.exists():
                logger.warning("AssemblyAgent: logo not found at %s — skipping slide %d", logo_path, slide_number)
                continue

            pptx_slide.shapes.add_picture(
                str(logo_path),
                Inches(lx), Inches(ly),
                Inches(lw), Inches(lh),
            )

        prs.save(pptx_path)
        logger.debug("AssemblyAgent: brand logos added to %s", pptx_path)

    def run(self, state: DeckState) -> dict:
        """Assemble the PPTX and return a state update with the file path.

        Args:
            state: Must have `slides`, `layout_specs`, and `selected_images`
                   all populated by the preceding agents.

        Returns:
            Dict with `pptx_path` (str) and updated status/logs.
        """
        s = get_settings()
        s.ensure_dirs()  # Make sure the output directory exists before writing

        # Build a safe filename: keep only alphanumeric chars and underscores.
        brief = state.deck_brief
        base = "deck"
        if brief and brief.title:
            base = "".join(c if c.isalnum() or c == "_" else "_" for c in brief.title)
            base = base[:50]  # Limit length to avoid OS path-length issues
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = str(Path(s.output_dir) / f"{base}_{timestamp}.pptx")

        # Convert the state's list fields into dicts keyed by slide_number so
        # PPTBuilder can look up a layout or image in O(1) per slide.
        layouts: Dict[int, LayoutSpec] = {
            spec.slide_number: spec for spec in state.layout_specs
        }
        images: Dict[int, str] = {
            si.slide_number: si.selected_image_path
            for si in state.selected_images
            if si.selected_image_path  # Skip slides where no image was selected
        }

        builder = PPTBuilder(output_path=output_path)
        builder.build(
            slides=state.slides,
            layouts=layouts,
            images=images,
        )

        # Add ML arteka logos per brand spec when brand is mobilelive / mlarteka.
        if state.brand.lower() in ("mobilelive", "mlarteka", "default"):
            self._add_brand_logos(output_path, state.slides, layouts)

        log_entry = f"AssemblyAgent: PPTX written → {output_path}"
        logger.info(log_entry)

        return {
            "pptx_path": output_path,
            "status": "assembly_complete",
            "execution_logs": state.execution_logs + [log_entry],
        }
